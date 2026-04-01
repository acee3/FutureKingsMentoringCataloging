from io import BytesIO

from pptx import Presentation


def extract_slide_text_from_pptx_bytes(pptx_bytes) -> list[str]:
    presentation = Presentation(BytesIO(pptx_bytes))
    slides = []

    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides.append("\n".join(slide_text))

    return slides
